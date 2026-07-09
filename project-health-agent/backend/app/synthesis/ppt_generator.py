"""
PPTX generator — builds an 8-slide executive presentation using python-pptx.

Slide structure:
1. Title slide — reporting period, portfolio name
2. Portfolio overview — RAG distribution (Doughnut/Pie chart) + KPI cards
3. Average Project Scores — Bar chart of metrics
4. Trend over time — Macro trends across the portfolio
5. RAG Status Shifts — Projects that changed health status
6. Top emerging risks — 2-column layout
7. Notable wins / green projects — 2-column layout
8. Executive recommendations
9. Appendix — Detailed per-project data table
"""

from __future__ import annotations

import logging
from datetime import datetime
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

# Modern Light Theme Colors
COLORS = {
    "green": RGBColor(16, 185, 129),
    "amber": RGBColor(245, 158, 11),
    "red": RGBColor(239, 68, 68),
    "bg": RGBColor(248, 250, 252),           # slate-50
    "card_bg": RGBColor(255, 255, 255),      # white
    "text_primary": RGBColor(15, 23, 42),    # slate-900
    "text_secondary": RGBColor(71, 85, 105), # slate-600
    "accent": RGBColor(37, 99, 235),         # blue-600
}

def _set_slide_background(slide, color: RGBColor = COLORS["bg"]) -> None:
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _add_title_text(slide, text: str, left: float = 0.5, top: float = 0.3, width: float = 9.0, height: float = 0.8, font_size: int = 28, color: RGBColor = COLORS["text_primary"], bold: bool = True) -> None:
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = PP_ALIGN.LEFT

def _add_body_text(slide, text: str, left: float = 0.5, top: float = 1.5, width: float = 9.0, height: float = 5.0, font_size: int = 14, color: RGBColor = COLORS["text_secondary"], bold: bool = False) -> None:
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(6)

def _add_bullet_list(slide, items: list[str], left: float = 0.5, top: float = 1.5, width: float = 9.0, height: float = 5.0, font_size: int = 14, color: RGBColor = COLORS["text_secondary"]) -> None:
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.level = 0

def _add_kpi_card(slide, title: str, value: str, left: float, top: float, width: float = 2.0, height: float = 1.0, value_color: RGBColor = COLORS["accent"]):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["card_bg"]
    shape.line.color.rgb = RGBColor(226, 232, 240) # slate-200 boundary

    tf = shape.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS["text_secondary"]
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = str(value)
    p2.font.size = Pt(24)
    p2.font.color.rgb = value_color
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER

def _add_styled_table(slide, headers: list[str], rows_data: list[list[str]], left: float = 0.5, top: float = 1.2, width: float = 9.0):
    rows_count = min(len(rows_data) + 1, 15)  # cap at 15
    cols_count = len(headers)
    table_shape = slide.shapes.add_table(rows_count, cols_count, Inches(left), Inches(top), Inches(width), Inches(rows_count * 0.4))
    table = table_shape.table

    # Headers
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["accent"]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.bold = True

    # Data
    for r in range(1, rows_count):
        row_data = rows_data[r-1]
        for c in range(cols_count):
            cell = table.cell(r, c)
            val = str(row_data[c]) if c < len(row_data) else ""
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["card_bg"]
            
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = COLORS["text_secondary"]
                
            # Color RAG status columns automatically
            val_lower = val.lower()
            if val_lower in ["green", "amber", "red"]:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS[val_lower]
                for paragraph in cell.text_frame.paragraphs:
                    # White text for red/green, very dark text for amber for contrast
                    paragraph.font.color.rgb = RGBColor(0,0,0) if val_lower == "amber" else RGBColor(255,255,255)
                    paragraph.font.bold = True

def generate_monthly_pptx(
    output_path: str,
    portfolio_name: str,
    period_start: datetime,
    period_end: datetime,
    project_snapshots: dict[str, list[dict[str, Any]]],
    synthesis_data: dict[str, Any],
    rag_counts: dict[str, int],
) -> str:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    total_projects = len(project_snapshots)

    # ── Slide 1: Title ──
    slide1 = prs.slides.add_slide(blank_layout)
    _set_slide_background(slide1)
    _add_title_text(slide1, portfolio_name, top=2.5, font_size=44, color=COLORS["accent"])
    _add_body_text(
        slide1,
        f"Executive Monthly Portfolio Review\n"
        f"{period_start.strftime('%B %d, %Y')} — {period_end.strftime('%B %d, %Y')}\n\n"
        f"Total Projects Analyzed: {total_projects}",
        top=4.0, font_size=18, color=COLORS["text_secondary"]
    )

    # ── Slide 2: Portfolio Overview Dashboard ──
    slide2 = prs.slides.add_slide(blank_layout)
    _set_slide_background(slide2)
    _add_title_text(slide2, "Portfolio Health Overview")
    
    exec_summary = synthesis_data.get("executive_summary", "")
    if exec_summary:
        _add_body_text(slide2, exec_summary, top=1.2, height=1.5, font_size=14, color=COLORS["text_primary"])

    # KPI Cards
    healthy_pct = (rag_counts.get("green", 0) / total_projects * 100) if total_projects else 0
    at_risk_pct = (rag_counts.get("red", 0) / total_projects * 100) if total_projects else 0
    
    _add_kpi_card(slide2, "Total Projects", str(total_projects), left=0.5, top=2.5)
    _add_kpi_card(slide2, "Healthy", f"{healthy_pct:.0f}%", left=2.8, top=2.5, value_color=COLORS["green"])
    _add_kpi_card(slide2, "At Risk", f"{at_risk_pct:.0f}%", left=5.1, top=2.5, value_color=COLORS["red"])

    # Pie Chart
    chart_data = ChartData()
    chart_data.categories = ["Green", "Amber", "Red"]
    chart_data.add_series("Projects", (rag_counts.get("green", 0), rag_counts.get("amber", 0), rag_counts.get("red", 0)))
    chart = slide2.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(1.5), Inches(3.8), Inches(6.0), Inches(3.2), chart_data).chart
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    
    try:
        colors = [COLORS["green"], COLORS["amber"], COLORS["red"]]
        series = chart.series[0]
        for i, color in enumerate(colors):
            point = series.points[i]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = color
    except Exception as e:
        logger.warning(f"Could not color pie chart: {e}")

    # ── Slide 3: Average Metrics Bar Chart ──
    slide3 = prs.slides.add_slide(blank_layout)
    _set_slide_background(slide3)
    _add_title_text(slide3, "Average Project Scores across Portfolio")
    
    # Calculate averages
    avg_scores = {"Schedule": 0, "Budget": 0, "Milestones": 0, "Blockers": 0, "Sentiment": 0}
    if total_projects > 0:
        for snaps in project_snapshots.values():
            if snaps:
                latest = snaps[-1]
                avg_scores["Schedule"] += latest.get("schedule_score") or 0
                avg_scores["Budget"] += latest.get("budget_score") or 0
                avg_scores["Milestones"] += latest.get("milestone_score") or 0
                avg_scores["Blockers"] += latest.get("blocker_score") or 0
                avg_scores["Sentiment"] += latest.get("sentiment_score") or 0
                
        for k in avg_scores:
            avg_scores[k] = avg_scores[k] / total_projects

    bar_chart_data = CategoryChartData()
    bar_chart_data.categories = list(avg_scores.keys())
    bar_chart_data.add_series("Average Score", tuple(avg_scores.values()))
    
    bchart = slide3.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.0), Inches(1.5), Inches(8.0), Inches(5.0), bar_chart_data).chart
    bchart.has_legend = False
    
    try:
        series = bchart.series[0]
        for i, pt in enumerate(series.points):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = COLORS["accent"]
    except Exception as e:
        pass


    # ── Slide 4: Cross-Project Trends ──
    slide4 = prs.slides.add_slide(blank_layout)
    _set_slide_background(slide4)
    _add_title_text(slide4, "Cross-Project Macro Trends")
    
    trends = synthesis_data.get("trends", [])
    if trends:
        headers = ["Trend Description", "Direction", "Severity", "Affected Projects"]
        rows = []
        for t in trends[:8]:
            rows.append([
                t.get("description", ""), 
                str(t.get("direction", "")).upper(), 
                str(t.get("severity", "")).upper(), 
                ", ".join(t.get("affected_projects", []))
            ])
        _add_styled_table(slide4, headers, rows)
    else:
        _add_body_text(slide4, "No significant macro trends identified.", top=2.0)

    # ── Slide 5: RAG Status Shifts ──
    slide5 = prs.slides.add_slide(blank_layout)
    _set_slide_background(slide5)
    _add_title_text(slide5, "RAG Status Shifts (Period-over-Period)")
    
    rag_shifts = synthesis_data.get("rag_shifts", [])
    if rag_shifts:
        headers = ["Project", "Previous Status", "Current Status", "Primary Reason"]
        rows = []
        for shift in rag_shifts[:8]:
            rows.append([
                shift.get("project", "Unknown"),
                str(shift.get("from", "N/A")).upper(),
                str(shift.get("to", "N/A")).upper(),
                shift.get("reason", "")
            ])
        _add_styled_table(slide5, headers, rows)
    else:
        _add_body_text(slide5, "No projects shifted RAG status during this period.", top=2.0)

    # ── Slide 6: Emerging Risks & Blockers ──
    slide6 = prs.slides.add_slide(blank_layout)
    _set_slide_background(slide6)
    _add_title_text(slide6, "Top Emerging Risks & Blockers", color=COLORS["red"])
    
    risks = synthesis_data.get("emerging_risks", [])
    if risks:
        headers = ["Risk Description", "Urgency", "Affected Projects"]
        rows = []
        for r in risks[:8]:
            rows.append([
                r.get("description", ""),
                str(r.get("urgency", "")).upper(),
                ", ".join(r.get("affected_projects", []))
            ])
        _add_styled_table(slide6, headers, rows)
    else:
        _add_body_text(slide6, "No critical emerging risks identified.", top=2.0)

    # ── Slide 7: Notable Wins & Green Projects ──
    slide7 = prs.slides.add_slide(blank_layout)
    _set_slide_background(slide7)
    _add_title_text(slide7, "Notable Wins & Positive Signals", color=COLORS["green"])
    
    wins = synthesis_data.get("notable_wins", [])
    if wins:
        headers = ["Win / Positive Signal", "Associated Projects"]
        rows = []
        for w in wins[:8]:
            rows.append([
                w.get("description", ""),
                ", ".join(w.get("projects", []))
            ])
        _add_styled_table(slide7, headers, rows)
    else:
        _add_body_text(slide7, "No notable wins highlighted for this period.", top=2.0)

    # ── Slide 8: Executive Recommendations ──
    slide8 = prs.slides.add_slide(blank_layout)
    _set_slide_background(slide8)
    _add_title_text(slide8, "Strategic Recommendations", color=COLORS["accent"])
    
    recs = synthesis_data.get("recommendations", [])
    if recs:
        _add_bullet_list(slide8, recs[:7], top=1.5, font_size=18, color=COLORS["text_primary"])
    else:
        _add_body_text(slide8, "No recommendations at this time.", top=2.0)

    # ── Slide 9: Appendix - Detailed Table ──
    slide9 = prs.slides.add_slide(blank_layout)
    _set_slide_background(slide9)
    _add_title_text(slide9, "Appendix: Detailed Portfolio Metrics", font_size=24)
    
    headers = ["Project Name", "RAG", "Overall", "Schedule", "Budget", "Blocker", "Sentiment"]
    rows = []
    for name, snaps in project_snapshots.items():
        if snaps:
            latest = snaps[-1]
            rows.append([
                name[:25],
                str(latest.get("rag_status", "N/A")).upper(),
                f"{latest.get('weighted_score') or 0:.0f}",
                f"{latest.get('schedule_score') or 0:.0f}",
                f"{latest.get('budget_score') or 0:.0f}",
                f"{latest.get('blocker_score') or 0:.0f}",
                f"{latest.get('sentiment_score') or 0:.0f}"
            ])
    _add_styled_table(slide9, headers, rows)

    # Save
    prs.save(output_path)
    logger.info("PPTX generated: %s", output_path)
    return output_path
